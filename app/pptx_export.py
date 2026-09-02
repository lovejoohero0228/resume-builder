"""포트폴리오를 가로(16:9) PPTX로 생성.

app/index.html 의 포트폴리오 상세 카드(디자인 포트폴리오)와 동일한 색·요소를 그대로 쓰되,
세로로 쌓이던 레이아웃을 한 슬라이드(가로) 안에 들어가도록 재배치한다.
프로젝트 1개 = 슬라이드 1장 (분리 슬라이드 없음).

레이아웃:
- 헤더(PROJECT 번호·제목·메타·칩·역할태그). 제목이 길어 2줄이 되면 아래 요소를 밀어낸다.
- 본문: 문제정의 → 나의역할 → 주요성과 를 읽기 순서 그대로 **균등 2컬럼에 높이 균형 분할**
  (한쪽만 넘치지 않게). 폰트는 **프로젝트별로** 콘텐츠 밀도에 맞는 최대 크기를 고른다.
- 다이어그램: 같은 슬라이드 **하단 풀폭 스트립**에 배치 (여러 장이면 가로로 나란히).

견고성: 텍스트 박스는 자동 축소(TEXT_TO_FIT)로 실제 PowerPoint 렌더 편차에서도 넘치지 않게 하고,
한글/영문 폭을 구분해 줄 수를 추정한다.

portfolio_pptx_bytes(identity, target, parts, lang, root_dir) -> (bytes, warnings)
"""
from __future__ import annotations
import io
import math
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ---- app/index.html 의 --pri 계열 팔레트와 동일 ----
PRI = RGBColor(0x1E, 0x3A, 0x8A)
PRI_INK = RGBColor(0x1B, 0x31, 0x68)
INK = RGBColor(0x1F, 0x29, 0x37)
BODY = RGBColor(0x37, 0x45, 0x5A)
MUT = RGBColor(0x6B, 0x74, 0x80)
TINT = RGBColor(0xEE, 0xF1, 0xFB)
CHIP = RGBColor(0xE0, 0xE6, 0xF6)
LINE = RGBColor(0xE5, 0xE8, 0xF1)
PILL_BORDER = RGBColor(0xD7, 0xDC, 0xF0)
IMPACT_BG = RGBColor(0xE9, 0xF0, 0xFF)
IMPACT_BORDER = RGBColor(0xD3, 0xDE, 0xF6)
IMPACT_TEXT = RGBColor(0x22, 0x34, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xED, 0xEF, 0xF4)
CARD_BORDER = RGBColor(0xE2, 0xE5, 0xEE)
FONT = "Malgun Gothic"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_W = Inches(SLIDE_W_IN)
SLIDE_H = Inches(SLIDE_H_IN)

CARD_MARGIN = 0.24     # 슬라이드 가장자리 -> 카드 가장자리
CARD_PAD_X = 0.3       # 카드 가장자리 -> 콘텐츠
CARD_PAD_TOP = 0.24
CARD_PAD_BOTTOM = 0.22
COL_GAP = 0.4
LEFT_RATIO = 0.615     # 왼쪽(문제 정의·나의 역할) : 오른쪽(주요 성과) 폭 비율

CHAR_W = 0.95        # pt당 문자 폭 추정 계수 (한글 기준, 넉넉하게 잡아 넘침 방지)
LINE_GAP = 1.3


# ---------- app/index.html 의 DIAGRAMS 매핑 파싱 ----------

def diagram_map(root_dir: str) -> dict:
    path = os.path.join(root_dir, "app", "index.html")
    try:
        with open(path, encoding="utf-8") as f:
            html = f.read()
    except OSError:
        return {}
    m = re.search(r"const\s+DIAGRAMS\s*=\s*\{(.*?)\}\s*;", html, re.S)
    if not m:
        return {}
    out = {}
    for em in re.finditer(r"(\w+)\s*:\s*\[(.*?)\]", m.group(1), re.S):
        files = [f.replace("\\'", "'") for f in re.findall(r"'((?:[^'\\]|\\.)*)'", em.group(2))]
        if files:
            out[em.group(1)] = files
    return out


# ---------- 텍스트 폭 추정 (자동 줄바꿈 높이 계산용) ----------

def _text_width_units(s):
    """문자열의 시각적 폭을 em 단위로 추정. 한중일(CJK)은 ~1.0em(정사각),
    영문·숫자·기호는 ~0.55em(좁음). 이래야 영/한 혼용에서 줄 수를 정확히 잡는다."""
    u = 0.0
    for ch in str(s):
        if ('가' <= ch <= '힣' or '぀' <= ch <= 'ヿ'
                or '一' <= ch <= '鿿' or '㄰' <= ch <= '㆏'):
            u += 1.0
        elif ch == ' ':
            u += 0.35
        else:
            u += 0.55
    return u


def est_lines(text, width_in, size_pt):
    if not text:
        return 0
    # 한 줄에 들어가는 em 폭 (CHAR_W 를 CJK 1em 폭 계수로 사용, 넉넉하게 잡아 넘침 방지)
    units_per_line = max(6.0, (width_in * 72) / (size_pt * CHAR_W))
    return max(1, math.ceil(_text_width_units(text) / units_per_line))


def est_h(text, width_in, size_pt, line_gap=LINE_GAP):
    return est_lines(text, width_in, size_pt) * (size_pt * line_gap) / 72.0


def _hurdle_lines(text):
    lines = [l.strip().lstrip("-").strip() for l in str(text or "").split("\n")]
    return [l for l in lines if l]


# ---------- 역할 → 태그(리드/참여) + 상세 (index.html roleScope/roleTagHtml 포팅) ----------

_TRIG = re.compile(r"리드|리더|주도|단독|참여|Lead|Sole|Drove|Contributor", re.I)


def role_scope(role, is_ko):
    if re.search(r"참여|contributor", role or "", re.I):
        return "프로젝트 참여" if is_ko else "Project Contributor"
    return "프로젝트 리더" if is_ko else "Project Lead"


def role_tag_parts(role, is_ko):
    if not role:
        return None
    scope = role_scope(role, is_ko)
    m = re.match(r"^([^(（]*)[(（]([^)）]*)\)?", role)
    pre = m.group(1) if m else role
    paren = m.group(2) if m else ""
    pre_rest = _TRIG.sub("", pre, count=1)
    pre_rest = re.sub(r"\s+", " ", pre_rest).strip(" ·,")
    detail = " · ".join(x for x in [pre_rest.strip(), paren.strip()] if x)
    return scope, detail


# ---------- 저수준 shape 헬퍼 (좌표는 전부 inch float) ----------

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line_color=None, line_w=0.75, radius=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if radius is not None:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def _textbox(slide, x, y, w, h, anchor=None, fit=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    if fit:
        # 안전망: 추정보다 텍스트가 길어도 박스 안에서 폰트를 살짝 줄여 겹침/넘침 방지
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tb, tf


def _run(p, text, size, color=INK, bold=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.name = FONT
    return r


def _rich(p, text, size, color=INK, bold=False):
    for i, seg in enumerate(str(text).split("**")):
        if seg:
            _run(p, seg, size, color, bold or (i % 2 == 1))


def _label(slide, x, y, w, h, text, size, color, bold=True):
    _, tf = _textbox(slide, x, y, w, h)
    _run(tf.paragraphs[0], text, size, color, bold)
    return tf


def _chip_width(text, size, pad=0.115):
    return max(0.36, len(text) * size * 0.0105 + pad * 2)


def _draw_chip(slide, x, y, text, fill, text_color, size=9.5, bold=True, h=0.27, border=None):
    w = _chip_width(text, size)
    sh = _rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
               line_color=border, line_w=0.75, radius=0.5)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, size, text_color, bold)
    return w


def _layout_chips(x0, max_w, chips, size=9.5, h=0.27, gap=0.09, line_gap=0.1):
    """flex-wrap 처럼 가로 배치 좌표를 계산 (그리기와 측정이 같은 값을 쓰도록 분리)."""
    x, rows = x0, 0
    placed = []
    for text, fill, color, border in chips:
        w = _chip_width(text, size)
        if x + w > x0 + max_w and x > x0:
            x = x0
            rows += 1
        placed.append((text, fill, color, border, x, rows))
        x += w + gap
    return placed, rows + 1


def _chip_rows_height(chips, max_w, size=9.5, h=0.27, line_gap=0.08):
    if not chips:
        return 0.0
    _, rows = _layout_chips(0.0, max_w, chips, size=size, h=h, line_gap=line_gap)
    return rows * h + (rows - 1) * line_gap


# ---------- 컨텐츠 블록 모델 (측정과 렌더가 같은 높이 계산을 공유) ----------
#
# `compact`(0~1)는 분량이 많은 프로젝트에서 슬라이드 안에 다 들어가도록 여백만 줄이는
# 값이다 (폰트 크기는 덱 전체에서 절대 바뀌지 않는다 — 이게 이번에 통일한 부분).

PILL_LABEL_H = 0.25
KH_H = 0.34


def _g(v, compact, factor=0.4):
    """여백/간격 상수를 compact 만큼 줄인다 (폰트 크기는 건드리지 않음)."""
    return v * (1 - factor * compact)


def _lg(compact):
    return LINE_GAP * (1 - 0.16 * compact)


def _block_height(block, width_in, size, compact=0.0):
    t = block["type"]
    if t == "kh":
        return _g(KH_H, compact, 0.18)
    if t == "pill_text":
        return _g(PILL_LABEL_H, compact, 0.18) + est_h(block["text"], width_in, size, _lg(compact)) + _g(0.11, compact)
    if t == "pill_bullets":
        h = _g(PILL_LABEL_H, compact, 0.18)
        for it in block["items"]:
            h += est_h(it, width_in - 0.18, size, _lg(compact)) + _g(0.065, compact)
        return h + _g(0.06, compact)
    if t == "group":
        h = 0.0
        if block["label"]:
            h += est_h(block["label"], width_in, size, _lg(compact)) + _g(0.07, compact)
        for it in block["items"]:
            h += est_h(it, width_in - 0.24, size - 0.5, _lg(compact)) + _g(0.06, compact)
        return h + _g(0.09, compact)
    if t == "impact":
        return est_h(block["text"], width_in - 0.58, size, _lg(compact)) + _g(0.22, compact, 0.2) + _g(0.11, compact)
    return 0.0


def _column_height(blocks, width_in, size, compact=0.0):
    return sum(_block_height(b, width_in, size, compact) for b in blocks)


def _draw_block(slide, block, x, y, width_in, size, compact=0.0):
    t = block["type"]
    if t == "kh":
        kh_h = _g(KH_H, compact, 0.18)
        _rect(slide, x, y + 0.035, 0.05, 0.17, PRI)
        _label(slide, x + 0.13, y - 0.02, width_in - 0.13, 0.25, block["text"], size + 1.5, PRI_INK, True)
        _rect(slide, x, y + kh_h - 0.035, width_in, Pt(1.1) / 914400, LINE)
        return kh_h

    if t == "pill_text":
        pill_h = _g(PILL_LABEL_H, compact, 0.18)
        _draw_chip(slide, x, y, block["label"], TINT, PRI, size=max(7.5, size - 1.5), bold=True,
                    h=PILL_LABEL_H - 0.02, border=PILL_BORDER)
        ty = y + pill_h
        th = est_h(block["text"], width_in, size, _lg(compact))
        _, tf = _textbox(slide, x, ty, width_in, th + 0.06, fit=True)
        tf.paragraphs[0].line_spacing = _lg(compact)
        _rich(tf.paragraphs[0], block["text"], size, BODY)
        return pill_h + th + _g(0.11, compact)

    if t == "pill_bullets":
        pill_h = _g(PILL_LABEL_H, compact, 0.18)
        _draw_chip(slide, x, y, block["label"], TINT, PRI, size=max(7.5, size - 1.5), bold=True,
                    h=PILL_LABEL_H - 0.02, border=PILL_BORDER)
        cy = y + pill_h
        for it in block["items"]:
            ih = est_h(it, width_in - 0.18, size, _lg(compact))
            _, tf = _textbox(slide, x, cy, width_in, ih + 0.05, fit=True)
            p = tf.paragraphs[0]
            p.line_spacing = _lg(compact)
            _run(p, "– ", size, PRI, True)
            _rich(p, it, size, BODY)
            cy += ih + _g(0.065, compact)
        return (cy - y) + _g(0.06, compact)

    if t == "group":
        cy = y
        if block["label"]:
            lh = est_h(block["label"], width_in, size, _lg(compact))
            _rect(slide, x, cy + 0.06, 0.08, 0.08, PRI)
            _, tf = _textbox(slide, x + 0.16, cy, width_in - 0.16, lh + 0.05)
            p = tf.paragraphs[0]
            p.line_spacing = _lg(compact)
            _run(p, block["label"], size, PRI_INK, True)
            cy += lh + _g(0.07, compact)
        for it in block["items"]:
            ih = est_h(it, width_in - 0.24, size - 0.5, _lg(compact))
            _, tf = _textbox(slide, x + 0.18, cy, width_in - 0.18, ih + 0.05, fit=True)
            p = tf.paragraphs[0]
            p.line_spacing = _lg(compact)
            _run(p, "▹ ", size - 0.5, PRI, True)
            _rich(p, it, size - 0.5, BODY)
            cy += ih + _g(0.06, compact)
        return (cy - y) + _g(0.09, compact)

    if t == "impact":
        th = est_h(block["text"], width_in - 0.58, size, _lg(compact))
        pad = _g(0.22, compact, 0.2)
        box_h = th + pad
        _rect(slide, x, y, width_in, box_h, IMPACT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              line_color=IMPACT_BORDER, line_w=0.75, radius=0.12)
        _, mtf = _textbox(slide, x + 0.14, y + pad / 2 - 0.04, 0.3, 0.3)
        _run(mtf.paragraphs[0], "✦", size + 1, PRI, True)
        _, tf = _textbox(slide, x + 0.44, y + pad / 2 - 0.03, width_in - 0.58, th + 0.06,
                         anchor=MSO_ANCHOR.MIDDLE, fit=True)
        p = tf.paragraphs[0]
        p.line_spacing = _lg(compact)
        _rich(p, block["text"], size, IMPACT_TEXT, True)
        return box_h + _g(0.11, compact)
    return 0.0


def _render_column(slide, blocks, x, y0, width_in, size, compact=0.0):
    y = y0
    for b in blocks:
        y += _draw_block(slide, b, x, y, width_in, size, compact)
    return y


# ---------- 프로젝트 → 블록 모델 ----------

def _project_block_list(part, lang):
    """프로젝트 본문을 읽기 순서(문제정의 → 나의역할 → 주요성과)의 단일 블록 리스트로.
    이 리스트를 두 컬럼에 높이 균형으로 나눠 배치한다."""
    is_ko = lang != "en"
    L = {"prob": "문제 정의", "role": "나의 역할", "imp": "주요 성과"} if is_ko \
        else {"prob": "PROBLEM", "role": "WHAT I DID", "imp": "IMPACT"}
    blocks = []

    problem = part.get("problem") or {}
    if problem.get("goal") or problem.get("hurdle"):
        blocks.append({"type": "kh", "text": L["prob"]})
        if problem.get("goal"):
            blocks.append({"type": "pill_text",
                           "label": "풀고자 한 문제" if is_ko else "Goal", "text": problem["goal"]})
        hlines = _hurdle_lines(problem.get("hurdle"))
        if hlines:
            blocks.append({"type": "pill_bullets",
                           "label": "제약과 어려움" if is_ko else "Constraints & hurdles", "items": hlines})

    role_groups = [g for g in (part.get("role_groups") or []) if not g.get("hidden")]
    if role_groups:
        blocks.append({"type": "kh", "text": L["role"]})
        for g in role_groups:
            blocks.append({"type": "group", "label": g.get("label") or "", "items": g.get("items") or []})

    impact = [it for it in (part.get("impact") or []) if not it.get("hidden")]
    if impact:
        blocks.append({"type": "kh", "text": L["imp"]})
        for it in impact:
            blocks.append({"type": "impact", "text": it.get("value") or ""})

    return blocks


def _split_columns(blocks, col_w, size, compact=0.0):
    """블록 리스트를 읽기 순서를 유지하며 두 컬럼으로 높이 균형 분할.
    컬럼 끝에 섹션 헤더(kh)만 남지 않게 보정한다."""
    if not blocks:
        return [], []
    hs = [_block_height(b, col_w, size, compact) for b in blocks]
    total = sum(hs)
    # 누적 높이가 절반을 처음 넘는 지점을 분할점으로
    acc, k = 0.0, len(blocks)
    for i, h in enumerate(hs):
        if acc + h > total / 2 and i > 0:
            # i 앞에서 자를지(acc) i 뒤에서 자를지(acc+h) 중 더 균형인 쪽
            k = i if abs(acc - total / 2) <= abs(acc + h - total / 2) else i + 1
            break
        acc += h
    k = max(1, min(k, len(blocks) - 1)) if len(blocks) > 1 else len(blocks)
    # 컬럼1의 마지막이 kh(헤더)면 다음 컬럼으로 넘겨 고아 헤더 방지
    while k > 1 and blocks[k - 1]["type"] == "kh":
        k -= 1
    return blocks[:k], blocks[k:]


_SIZE_STEPS = (11.5, 11, 10.5, 10, 9.5, 9, 8.5, 8, 7.5, 7, 6.5)
_COMPACT_STEPS = (0.0, 0.25, 0.5, 0.75, 1.0)


SAFETY_MARGIN_IN = 0.16  # 줄바꿈 추정 오차 대비 여유 (실제 PowerPoint 렌더 편차 흡수)


def _fits(blocks, col_w, text_h, size, compact=0.0):
    budget = text_h - SAFETY_MARGIN_IN
    c1, c2 = _split_columns(blocks, col_w, size, compact)
    return (_column_height(c1, col_w, size, compact) <= budget and
            _column_height(c2, col_w, size, compact) <= budget)


def _min_compact(blocks, col_w, text_h, size):
    """이 폰트 크기로 두 컬럼에 다 들어가게 할 최소 압축값 (0=넉넉, 1=최대). 안 되면 None."""
    for c in _COMPACT_STEPS:
        if _fits(blocks, col_w, text_h, size, c):
            return c
    return None


def _pick_project_layout(blocks, col_w, text_h):
    """프로젝트별로 두 컬럼에 들어가는 '가장 큰' 폰트 크기 + 그때의 최소 압축값을 고른다.
    (덱 통일 폰트 대신 슬라이드별 최적화 — 분량 가벼운 프로젝트는 크게, 무거운 것만 작게.)"""
    for size in _SIZE_STEPS:
        c = _min_compact(blocks, col_w, text_h, size)
        if c is not None:
            return size, c
    return _SIZE_STEPS[-1], 1.0


# ---------- 카드 프레임 (배경 + 흰 카드) ----------

def _card_frame(slide):
    """페이지 배경 + 흰 카드를 그리고, 카드 안쪽 콘텐츠 가능 영역(x0,y0,x1,y1)을 inch로 반환."""
    _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, PAGE_BG)
    cw = SLIDE_W_IN - CARD_MARGIN * 2
    ch = SLIDE_H_IN - CARD_MARGIN * 2
    _rect(slide, CARD_MARGIN, CARD_MARGIN, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          line_color=CARD_BORDER, line_w=1, radius=0.035)
    x0 = CARD_MARGIN + CARD_PAD_X
    y0 = CARD_MARGIN + CARD_PAD_TOP
    x1 = SLIDE_W_IN - CARD_MARGIN - CARD_PAD_X
    y1 = SLIDE_H_IN - CARD_MARGIN - CARD_PAD_BOTTOM
    return x0, y0, x1, y1


# ---------- 헤더 (레이아웃 계산과 그리기를 분리해 모든 슬라이드가 같은 기준선을 쓰게 한다) ----------

TITLE_PT = 20


def _header_layout(part, lang, content_w):
    """헤더 요소들의 y 오프셋(카드 콘텐츠 top=y0 기준)을 계산. 제목 줄 수를 실제로
    측정해 여러 줄 제목이면 메타·칩·본문을 그만큼 아래로 민다 (그리기·측정 공용)."""
    is_ko = lang != "en"
    title = part.get("title") or ""
    # 제목 줄 수를 넉넉히(1.06배) 잡아 실제 PowerPoint 에서 2줄이 돼도 겹치지 않게
    title_lines = max(1, math.ceil(est_lines(title, content_w, TITLE_PT) * 1.06))
    title_h = title_lines * (TITLE_PT * 1.18) / 72.0
    y_title = 0.21
    y_meta = y_title + title_h + 0.05
    meta = "   ·   ".join(x for x in [part.get("org"), part.get("period")] if x)
    y_chips = y_meta + (0.24 if meta else 0.0)
    chips = [(a, PRI, WHITE, None) for a in (part.get("angles") or [])] + \
            [(t, CHIP, PRI_INK, None) for t in (part.get("tags") or [])]
    chip_h = _chip_rows_height(chips, content_w) if chips else 0.0
    div_y = y_chips + chip_h + 0.08
    role = role_tag_parts(part.get("role") or "", is_ko)
    row_y = div_y + 0.1
    row_h = 0.28 if role else 0.0
    body_top = div_y + 0.1 + row_h + (0.08 if role else 0.02)
    return {"title": title, "title_h": title_h, "y_title": y_title, "meta": meta,
            "y_meta": y_meta, "chips": chips, "y_chips": y_chips, "div_y": div_y,
            "role": role, "row_y": row_y, "body_top": body_top}


def _draw_header(slide, part, num, lang, x0, y0, x1, hl):
    content_w = x1 - x0
    _label(slide, x0, y0, content_w, 0.2, f"PROJECT {num:02d}", 10.5, PRI, True)
    _, ttf = _textbox(slide, x0, y0 + hl["y_title"], content_w, hl["title_h"] + 0.12)
    _run(ttf.paragraphs[0], hl["title"], TITLE_PT, INK, True)

    if hl["meta"]:
        _label(slide, x0, y0 + hl["y_meta"], content_w, 0.2, hl["meta"], 10.5, MUT, False)

    if hl["chips"]:
        placed, _rows = _layout_chips(x0, content_w, hl["chips"])
        for text, fill, color, border, cx, row in placed:
            _draw_chip(slide, cx, y0 + hl["y_chips"] + row * (0.27 + 0.08), text, fill, color, border=border)

    _rect(slide, x0, y0 + hl["div_y"], content_w, Pt(1) / 914400, LINE)

    if hl["role"]:
        scope, detail = hl["role"]
        ry = y0 + hl["row_y"]
        w = _draw_chip(slide, x0, ry, scope, PRI, WHITE, size=10, bold=True, h=0.28)
        if detail:
            _label(slide, x0 + w + 0.14, ry + 0.035, content_w - w - 0.14, 0.24, detail, 10.5, MUT, False)


# ---------- 슬라이드 조립 ----------

# 다이어그램 하단 스트립: 이미지 장수에 따른 높이 (본문 텍스트 영역에서 미리 빼둔다)
def _img_strip_h(n):
    if n <= 0:
        return 0.0
    return 1.5 if n == 1 else (1.65 if n == 2 else 1.8)


def _draw_image_strip(slide, files, root_dir, lang, x0, w_total, y_top, band_h):
    """다이어그램들을 하단 풀폭 스트립에 가로로 나란히 배치 (비율 유지·중앙 정렬·테두리)."""
    is_ko = lang != "en"
    warnings, valid = [], []
    for fn in files:
        fp = os.path.join(root_dir, "assets", "diagrams", fn)
        if os.path.isfile(fp):
            valid.append(fp)
        else:
            warnings.append(f"이미지 없음: {fn}")
    if not valid:
        return warnings

    _label(slide, x0, y_top, w_total, 0.16, "다이어그램" if is_ko else "Diagram", 8.5, MUT, True)
    y_top += 0.2
    band_h -= 0.2
    if band_h < 0.3:
        return warnings

    n = len(valid)
    hgap = 0.16
    slot_w = (w_total - hgap * (n - 1)) / n
    for i, fp in enumerate(valid):
        sx = x0 + i * (slot_w + hgap)
        try:
            pic = slide.shapes.add_picture(fp, Emu(0), Emu(0))
        except Exception:
            warnings.append(f"이미지 삽입 실패: {os.path.basename(fp)}")
            continue
        iw, ih = pic.width, pic.height
        scale = min(Inches(slot_w) / iw, Inches(band_h) / ih)
        nw, nh = int(iw * scale), int(ih * scale)
        pic.width, pic.height = nw, nh
        pic.left = int(Inches(sx) + (Inches(slot_w) - nw) / 2)
        pic.top = int(Inches(y_top) + (Inches(band_h) - nh) / 2)
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
        pic.shadow.inherit = False
    return warnings


def add_project_slide(prs, part, num, lang, size, compact, col_w, gap, files, root_dir):
    slide = _blank_slide(prs)
    x0, y0, x1, y1 = _card_frame(slide)
    hl = _header_layout(part, lang, x1 - x0)
    _draw_header(slide, part, num, lang, x0, y0, x1, hl)
    body_top = y0 + hl["body_top"]

    strip_h = _img_strip_h(len(files))
    text_bottom = y1 - (strip_h + 0.14 if strip_h else 0.0)

    blocks = _project_block_list(part, lang)
    c1, c2 = _split_columns(blocks, col_w, size, compact)
    _render_column(slide, c1, x0, body_top, col_w, size, compact)
    _render_column(slide, c2, x0 + col_w + gap, body_top, col_w, size, compact)

    # 다이어그램: 하단 풀폭 스트립 (텍스트 영역과 이미 분리해 두어 겹치지 않음)
    warnings = []
    if files and strip_h:
        warnings += _draw_image_strip(slide, files, root_dir, lang, x0, x1 - x0,
                                      text_bottom + 0.14, strip_h)
    return warnings


def add_cover_slide(prs, identity, target, lang):
    is_ko = lang != "en"
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, PRI)
    _, tf = _textbox(slide, 1, 2.6, SLIDE_W_IN - 2, 1.6, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _run(p, target or ("포트폴리오" if is_ko else "Portfolio"), 15, RGBColor(0x9D, 0xC0, 0xF0), True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    _run(p2, (identity or {}).get("name") or "", 34, WHITE, True)
    if (identity or {}).get("tagline"):
        p3 = tf.add_paragraph()
        p3.space_before = Pt(10)
        _run(p3, identity["tagline"], 15, RGBColor(0xCF, 0xDD, 0xF2), False)


def portfolio_pptx_bytes(identity: dict, target: str, parts: list, lang: str, root_dir: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover_slide(prs, identity, target, lang)

    projects = [p for p in parts if p.get("kind") == "project"]

    # 콘텐츠 영역: 균등 2컬럼 (본문 블록을 두 컬럼에 균형 분할)
    content_w = SLIDE_W_IN - CARD_MARGIN * 2 - CARD_PAD_X * 2
    col_w = round((content_w - COL_GAP) / 2, 3)
    y1 = SLIDE_H_IN - CARD_MARGIN - CARD_PAD_BOTTOM

    dmap = diagram_map(root_dir)

    # 프로젝트별로 최적 폰트 크기·압축을 계산(이미지 있으면 하단 스트립만큼 텍스트 영역을 뺀다)
    # 후 렌더. 슬라이드마다 콘텐츠 밀도에 맞는 크기를 써서 가벼운 프로젝트는 크게 보인다.
    warnings = []
    for num, part in enumerate(projects, start=1):
        files = dmap.get(part.get("id"), [])
        blocks = _project_block_list(part, lang)
        body_top = _header_layout(part, lang, content_w)["body_top"]
        body_h = y1 - (CARD_MARGIN + CARD_PAD_TOP + body_top)
        strip_h = _img_strip_h(len(files))
        text_h = body_h - (strip_h + 0.14 if strip_h else 0.0)
        size, compact = _pick_project_layout(blocks, col_w, text_h)
        warnings += add_project_slide(prs, part, num, lang, size, compact,
                                      col_w, COL_GAP, files, root_dir)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue(), warnings
